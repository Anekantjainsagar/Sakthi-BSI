#!/usr/bin/env python3
"""
ENHANCED INTELLIGENT PROFESSIONAL AI SECURITY REPORT GENERATOR
With Risk Score, Industry Benchmark, Riskometer, and Better Text Wrapping
"""

import os
from gemini_config import call_gemini as _gemini_call, GEMINI_MODEL, GEMINI_API_KEYS
import json
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.dml import MSO_THEME_COLOR
import math


class EnhancedIntelligentReportGenerator:
    """
    Enhanced AI report generator with industry benchmarking,
    risk scoring, riskometer visualization, and improved text handling
    """
    
    COLORS = {
        'critical': RGBColor(211, 47, 47),
        'high': RGBColor(255, 111, 0),
        'medium': RGBColor(255, 193, 7),
        'low': RGBColor(76, 175, 80),
        'info': RGBColor(33, 150, 243),
        'primary': RGBColor(26, 35, 126),
        'secondary': RGBColor(69, 90, 100),
        'accent': RGBColor(0, 188, 212),
        'white': RGBColor(255, 255, 255),
        'light_grey': RGBColor(245, 245, 245),
        'dark_grey': RGBColor(66, 66, 66),
        'text': RGBColor(33, 33, 33),
        'green': RGBColor(76, 175, 80),
        'red': RGBColor(244, 67, 54),
    }
    
    # Industry benchmark ranges by sector
    INDUSTRY_BENCHMARKS = {
        'Technology': {'low': 72, 'high': 85},
        'Healthcare': {'low': 68, 'high': 82},
        'Finance': {'low': 75, 'high': 88},
        'Retail': {'low': 65, 'high': 78},
        'Manufacturing': {'low': 62, 'high': 75},
        'Education': {'low': 58, 'high': 72},
        'Government': {'low': 70, 'high': 83},
        'Default': {'low': 65, 'high': 78}
    }
    
    def __init__(self):
        self.use_ai = len(GEMINI_API_KEYS) > 0
        if self.use_ai:
            print(f"✅ Enhanced AI Report Generator initialized ({GEMINI_MODEL}, {len(GEMINI_API_KEYS)} keys)")
        else:
            print("⚠️ AI Report Generator: no Gemini keys available")
    
    def calculate_risk_score(self, phase4_data, phase5_data):
        """
        Calculate risk score (0-100) where LOWER is WORSE (more risk = lower score).
        Starts at 100 (perfect) and deducts points for each vulnerability found.
        """
        try:
            cves = phase4_data.get('cves_all', [])
            if not isinstance(cves, list): cves = []

            critical = len([c for c in cves if isinstance(c, dict) and c.get('cvss', 0) >= 9.0])
            high     = len([c for c in cves if isinstance(c, dict) and 7.0 <= c.get('cvss', 0) < 9.0])
            medium   = len([c for c in cves if isinstance(c, dict) and 4.0 <= c.get('cvss', 0) < 7.0])

            # Start at 100 (perfect security) and deduct for each vulnerability
            score = 100
            score -= critical * 15   # Critical CVE = -15 pts each
            score -= high     * 7    # High CVE    = -7  pts each
            score -= medium   * 2    # Medium CVE  = -2  pts each

            # Deduct for overall risk rating from Phase 5
            risk_rating = self._extract_nested(phase5_data, ['multidimensional_score', 'risk_rating'], 'Medium Risk')
            if 'Extreme' in str(risk_rating) or 'CRITICAL' in str(risk_rating).upper():
                score -= 25
            elif 'High' in str(risk_rating):
                score -= 15
            elif 'Medium' in str(risk_rating):
                score -= 8

            score = min(100, max(0, score))
            return int(score)
        except Exception as e:
            print(f"Error calculating risk score: {e}")
            return 35  # Default: poor security posture
    
    def get_industry_benchmark(self, industry, phase4_data):
        """
        Calculate industry benchmark as the percentage of serious vulnerabilities.
        Formula: ((Critical + High + Medium) / Grand Total of All CVEs) * 100
        Higher % = more of the vulnerabilities are serious = worse posture.
        Falls back to hardcoded sector average if no CVE data available.
        """
        cves = phase4_data.get('cves_all', [])
        if not isinstance(cves, list):
            cves = []

        total = len(cves)
        if total == 0:
            # No CVE data — fall back to hardcoded sector average
            benchmark_range = self.INDUSTRY_BENCHMARKS.get(industry, self.INDUSTRY_BENCHMARKS['Default'])
            return (benchmark_range['low'] + benchmark_range['high']) // 2

        critical = len([c for c in cves if isinstance(c, dict) and c.get('cvss', 0) >= 9.0])
        high     = len([c for c in cves if isinstance(c, dict) and 7.0 <= c.get('cvss', 0) < 9.0])
        medium   = len([c for c in cves if isinstance(c, dict) and 4.0 <= c.get('cvss', 0) < 7.0])

        benchmark = round(((critical + high + medium) / total) * 100)
        return benchmark
    
    def comprehensive_ai_analysis(self, phase1, phase2, phase3, phase4, phase5):
        """
        AI COMPREHENSIVELY ANALYZES ALL 5 PHASES with enhanced intelligence
        """
        print("\n🧠 AI performing comprehensive 5-phase analysis...")
        
        # Extract ALL data from all phases
        comprehensive_data = {
            "phase1_business": {
                "domain": phase1.get('domain') or phase1.get('url', '').replace('https://', '').replace('http://', '').split('/')[0] or 'Unknown',
                "company_name": phase1.get('company_name') or phase1.get('domain', 'Unknown').split('.')[0].title(),
                "industry": self._extract_nested(phase1, ['ai_analysis', 'company_overview', 'industry_vertical'], 'Technology'),
                "business_model": self._extract_nested(phase1, ['ai_analysis', 'company_overview', 'business_model'], 'B2B'),
                "revenue_model": self._extract_nested(phase1, ['ai_analysis', 'company_overview', 'revenue_model'], 'Unknown'),
            },
            "phase2_infrastructure": {
                "total_ips": len(phase2.get('ip_addresses', [])),
                "subdomains": len(phase2.get('subdomains', [])),
                "subdomain_list": phase2.get('subdomains', [])[:10],
                "open_ports": sum(len(ports) for ports in phase2.get('open_ports', {}).values()),
                "port_details": phase2.get('open_ports', {}),
                "ssl_issues": phase2.get('ssl_issues', []),
                "dns_records": phase2.get('dns_records', {}),
            },
            "phase3_application": {
                "technologies": phase3.get('1_application_discovery', {}).get('technologies', [])[:15],
                "cms": phase3.get('2_web_server_stack', {}).get('cms', ['Unknown'])[0] if isinstance(phase3.get('2_web_server_stack', {}).get('cms', []), list) else phase3.get('2_web_server_stack', {}).get('cms', 'Unknown'),
                "cms_version": phase3.get('2_web_server_stack', {}).get('cms_version', 'Unknown'),
                "javascript_libs": phase3.get('3_javascript_libraries', [])[:10],
                "analytics": phase3.get('4_analytics_tracking', []),
                "security_headers": phase3.get('7_security_posture', {}).get('security_headers', {}),
                "waf_detected": phase3.get('7_security_posture', {}).get('waf', {}).get('detected', False),
                "missing_headers": sum(1 for h in phase3.get('7_security_posture', {}).get('security_headers', {}).values() if not h.get('present')),
            },
            "phase4_vulnerabilities": {
                "total_cves": len(phase4.get('cves_all', [])),
                "critical_cves": len([c for c in phase4.get('cves_all', []) if isinstance(c, dict) and c.get('cvss', 0) >= 9.0]),
                "high_cves": len([c for c in phase4.get('cves_all', []) if isinstance(c, dict) and 7.0 <= c.get('cvss', 0) < 9.0]),
                "medium_cves": len([c for c in phase4.get('cves_all', []) if isinstance(c, dict) and 4.0 <= c.get('cvss', 0) < 7.0]),
                "low_cves": len([c for c in phase4.get('cves_all', []) if isinstance(c, dict) and c.get('cvss', 0) < 4.0]),
                "top_cves": sorted(
                    [c for c in phase4.get('cves_all', []) if isinstance(c, dict)],
                    key=lambda x: x.get('cvss', 0),
                    reverse=True
                )[:5],
                "cve_details": [
                    {
                        "id": c.get('cve', c.get('cve_id', 'Unknown')),
                        "cvss": c.get('cvss', 0),
                        "description": c.get('description', c.get('desc', 'No description'))[:200],
                        "technology": c.get('tech', c.get('product', 'Unknown'))
                    }
                    for c in sorted(
                        [c for c in phase4.get('cves_all', []) if isinstance(c, dict)],
                        key=lambda x: x.get('cvss', 0),
                        reverse=True
                    )[:5]
                ],
            },
            "phase5_risk": {
                "overall_score": self._extract_nested(phase5, ['multidimensional_score', 'overall_risk_score'], 5.0),
                "risk_rating": self._extract_nested(phase5, ['multidimensional_score', 'risk_rating'], 'MEDIUM'),
                "business_risk_level": self._extract_nested(phase5, ['business_risk', 'risk_level'], 'Medium'),
                "infra_risk_level": self._extract_nested(phase5, ['infrastructure_risk', 'risk_level'], 'Medium'),
                "app_risk_level": self._extract_nested(phase5, ['application_risk', 'risk_level'], 'Medium'),
                "business_impact_score": self._extract_nested(phase5, ['business_risk', 'business_impact_score'], 5.0),
            }
        }
        
        # Calculate risk score and benchmark
        risk_score = self.calculate_risk_score(phase4, phase5)
        industry = comprehensive_data['phase1_business']['industry']
        industry_benchmark = self.get_industry_benchmark(industry, phase4)
        
        comprehensive_data['risk_metrics'] = {
            'risk_score': risk_score,
            'industry_benchmark': industry_benchmark,
            'industry': industry
        }
        
        # AI COMPREHENSIVE ANALYSIS PROMPT
        prompt = f"""
You are a senior cybersecurity consultant analyzing a complete security assessment.
Perform COMPREHENSIVE, INTELLIGENT analysis of ALL 5 phases.

COMPLETE DATA FROM ALL 5 PHASES:
{json.dumps(comprehensive_data, indent=2)}

CONTEXT:
- Company Security Score: {risk_score}/100 (LOWER = worse security posture, HIGHER = better)
- Industry Benchmark: {industry_benchmark}/100 (average for this sector)
- Industry: {industry}

TASK: Create INTELLIGENT, CONSULTING-QUALITY analysis showing deep understanding.

REQUIREMENTS:
1. ANALYZE relationships between phases intelligently
2. IDENTIFY specific attack vectors and entry points from the data
3. PROVIDE business impact context (not just technical)
4. CREATE actionable, prioritized recommendations
5. SYNTHESIZE findings across all phases
6. Focus on PRACTICAL security improvements
7. Provide 3-5 specific "Cyberattack Entry Points" based on actual vulnerabilities found

OUTPUT STRICT JSON:
{{
  "comprehensive_analysis": {{
    "executive_summary": "2-3 sentences summarizing overall security posture and key risks",
    "critical_insight": "1 sentence most important finding",
    "attack_surface_analysis": "2 sentences analyzing exposure + vulnerabilities",
    "technology_risk_correlation": "2 sentences connecting tech stack to risks"
  }},
  "cyberattack_entry_points": [
    {{
      "title": "Specific vulnerability or weakness name (8-12 words max)",
      "description": "Concise explanation of the risk and potential impact (max 180 chars)",
      "severity": "Critical/High/Medium/Low"
    }},
    {{
      "title": "Another specific entry point",
      "description": "Brief, focused description of the attack vector",
      "severity": "Critical/High/Medium/Low"
    }},
    {{
      "title": "Third entry point",
      "description": "Clear explanation of exploitation risk",
      "severity": "Critical/High/Medium/Low"
    }}
  ],
  "intelligent_findings": {{
    "finding_1": {{
      "title": "Cross-Phase Critical Finding",
      "synthesis": "2 sentences connecting multiple phases",
      "business_impact": "1 sentence on business consequences",
      "severity": "CRITICAL/HIGH/MEDIUM/LOW",
      "affected_phases": ["phase names"]
    }},
    "finding_2": {{
      "title": "Strategic Security Gap",
      "synthesis": "2 sentences of intelligent analysis",
      "business_impact": "1 sentence on business impact",
      "severity": "CRITICAL/HIGH/MEDIUM/LOW",
      "affected_phases": ["phase names"]
    }},
    "finding_3": {{
      "title": "Infrastructure-Vulnerability Correlation",
      "synthesis": "2 sentences connecting infrastructure to vulnerabilities",
      "business_impact": "1 sentence on business consequences",
      "severity": "CRITICAL/HIGH/MEDIUM/LOW",
      "affected_phases": ["phase names"]
    }}
  }},
  "strategic_recommendations": {{
    "immediate_priorities": [
      "Priority 1: Specific actionable step (focus on highest risk)",
      "Priority 2: Another concrete action addressing critical gaps",
      "Priority 3: Third priority based on comprehensive analysis"
    ],
    "strategic_initiatives": [
      "Initiative 1: Long-term security program improvement",
      "Initiative 2: Strategic capability enhancement"
    ]
  }},
  "executive_summary": {{
    "headline": "One compelling sentence summarizing security posture",
    "risk_narrative": "2-3 sentences explaining the risk landscape and business impact",
    "key_metrics": {{
      "total_vulnerabilities": {comprehensive_data['phase4_vulnerabilities']['total_cves']},
      "critical_findings": {comprehensive_data['phase4_vulnerabilities']['critical_cves']},
      "high_findings": {comprehensive_data['phase4_vulnerabilities']['high_cves']}
    }}
  }}
}}

IMPORTANT: 
- Keep all text concise and scannable
- Entry point titles must be specific (not generic)
- Descriptions must fit in small boxes (180 chars max)
- Focus on ACTUAL findings from the data
- Be professional but direct
"""
        
        if self.use_ai:
            try:
                raw = _gemini_call(prompt, max_tokens=8192, temperature=0.3)
                if raw:
                    ai_analysis = self._safe_parse_json(raw)
                    if ai_analysis:
                        print("✅ AI comprehensive analysis complete")
                        return ai_analysis, comprehensive_data
                    else:
                        print("⚠️ AI returned unparseable response, using template")
            except Exception as e:
                print(f"⚠️ AI analysis failed: {e}, using template")
        return self._fallback_analysis(comprehensive_data), comprehensive_data

    def _safe_parse_json(self, raw: str) -> dict:
        """Safely parse JSON from AI response — handles fences, extra text, truncation"""
        if not raw:
            return {}
        text = raw.strip()
        # Strip markdown fences
        if '```json' in text:
            text = text.split('```json')[1].split('```')[0].strip()
        elif '```' in text:
            text = text.split('```')[1].split('```')[0].strip()
        # Try direct parse
        try:
            return json.loads(text)
        except Exception:
            pass
        # Try extracting first JSON object (handles extra text around JSON)
        import re
        match = re.search(r'\{[\s\S]*\}', text)
        if match:
            try:
                return json.loads(match.group(0))
            except Exception:
                pass
        return {}
    
    def _fallback_analysis(self, data):
        """Fallback template when AI is unavailable"""
        return {
            "comprehensive_analysis": {
                "executive_summary": f"Security assessment identified {data['phase4_vulnerabilities']['total_cves']} vulnerabilities across the technology stack.",
                "critical_insight": f"Found {data['phase4_vulnerabilities']['critical_cves']} critical vulnerabilities requiring immediate attention.",
                "attack_surface_analysis": "Infrastructure analysis reveals potential attack vectors through exposed services and outdated technologies.",
                "technology_risk_correlation": "Technology stack includes components with known vulnerabilities that increase organizational risk."
            },
            "cyberattack_entry_points": [
                {
                    "title": "Outdated Software Components with Known Vulnerabilities",
                    "description": "Multiple technology components running outdated versions expose the organization to known exploits.",
                    "severity": "High"
                },
                {
                    "title": "Missing Security Headers on Web Application",
                    "description": "Absence of critical security headers increases vulnerability to XSS and clickjacking attacks.",
                    "severity": "Medium"
                },
                {
                    "title": "Exposed Services on Public Infrastructure",
                    "description": "Services exposed to the internet provide potential entry points for unauthorized access.",
                    "severity": "High"
                }
            ],
            "intelligent_findings": {
                "finding_1": {
                    "title": "Critical Vulnerability Exposure",
                    "synthesis": f"Analysis identified {data['phase4_vulnerabilities']['critical_cves']} critical CVEs in production systems.",
                    "business_impact": "Could lead to data breaches and service disruption.",
                    "severity": "CRITICAL",
                    "affected_phases": ["Phase 4"]
                },
                "finding_2": {
                    "title": "Infrastructure Security Gaps",
                    "synthesis": "Exposed services and open ports create additional attack surface.",
                    "business_impact": "Increases risk of unauthorized access and reconnaissance.",
                    "severity": "HIGH",
                    "affected_phases": ["Phase 2"]
                },
                "finding_3": {
                    "title": "Application Security Concerns",
                    "synthesis": f"{data['phase3_application']['missing_headers']} security headers missing from web applications.",
                    "business_impact": "Exposes users to client-side attacks.",
                    "severity": "MEDIUM",
                    "affected_phases": ["Phase 3"]
                }
            },
            "strategic_recommendations": {
                "immediate_priorities": [
                    "Patch critical and high severity vulnerabilities within 30 days",
                    "Implement missing security headers on all web applications",
                    "Review and minimize exposed services on public infrastructure"
                ],
                "strategic_initiatives": [
                    "Establish continuous vulnerability management program",
                    "Implement security monitoring and threat detection capabilities"
                ]
            },
            "executive_summary": {
                "headline": "Security assessment reveals actionable improvements needed.",
                "risk_narrative": "The organization faces moderate security risks that can be addressed through systematic remediation efforts.",
                "key_metrics": {
                    "total_vulnerabilities": data['phase4_vulnerabilities']['total_cves'],
                    "critical_findings": data['phase4_vulnerabilities']['critical_cves'],
                    "high_findings": data['phase4_vulnerabilities']['high_cves']
                }
            }
        }
    
    def create_intelligent_presentation(self, analysis, data, output_path):
        """
        Create professional presentation with enhanced features
        """
        print("\n📊 Creating enhanced professional presentation...")
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Create slides
        self._create_cover_slide(prs, analysis, data)
        self._create_executive_summary_enhanced(prs, analysis, data)
        self._create_vulnerability_analysis(prs, analysis, data)
        self._create_findings_slide(prs, analysis, data)
        self._create_roadmap(prs, analysis, data)
        
        # Save
        domain = data['phase1_business']['domain']
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"Enhanced_SecurityReport_{domain}_{timestamp}.pptx"
        filepath = os.path.join(output_path, filename)
        
        prs.save(filepath)
        print(f"✅ Enhanced presentation saved: {filename}")
        
        return filepath
    
    def _create_cover_slide(self, prs, analysis, data):
        """Enhanced cover slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Background
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(7.5)
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.COLORS['primary']
        background.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
        tf = title_box.text_frame
        tf.word_wrap = True
        
        company = data['phase1_business']['company_name']
        
        p = tf.paragraphs[0]
        p.text = "Cybersecurity"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        
        p = tf.add_paragraph()
        p.text = "Intelligence Report"
        p.font.size = Pt(38)
        p.font.color.rgb = self.COLORS['accent']
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(10)
        
        # Company name
        company_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
        tf = company_box.text_frame
        p = tf.paragraphs[0]
        p.text = company
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        
        # Date
        date_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Generated: {datetime.now().strftime('%B %d, %Y')}"
        p.font.size = Pt(14)
        p.font.color.rgb = self.COLORS['accent']
        p.alignment = PP_ALIGN.CENTER
    
    def _add_wrapped_text(self, slide, text, left, top, width, height, font_size=11, bold=False, color=None):
        """Helper to add text with proper wrapping"""
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = 1  # Middle
        
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        if color:
            p.font.color.rgb = color
        else:
            p.font.color.rgb = self.COLORS['text']
        p.line_spacing = 1.2
        
        return textbox
    
    def _create_riskometer(self, slide, risk_score, left, top, size=2.5):
        """Create a riskometer gauge visualization"""
        # Background circle
        gauge_bg = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            left, top, Inches(size), Inches(size)
        )
        gauge_bg.fill.solid()
        gauge_bg.fill.fore_color.rgb = self.COLORS['light_grey']
        gauge_bg.line.color.rgb = self.COLORS['dark_grey']
        gauge_bg.line.width = Pt(2)
        
        # Lower score = worse security. Colour reflects risk severity correctly.
        if risk_score >= 75:
            needle_color = self.COLORS['green']    # Good security posture
            status = "Good"
        elif risk_score >= 55:
            needle_color = self.COLORS['medium']   # Medium risk
            status = "Medium"
        elif risk_score >= 35:
            needle_color = self.COLORS['high']     # High risk
            status = "High"
        else:
            needle_color = self.COLORS['critical'] # Critical — very poor posture
            status = "Critical"
        
        # Center circle with score
        center_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            left + Inches(size/3), top + Inches(size/3), 
            Inches(size/3), Inches(size/3)
        )
        center_circle.fill.solid()
        center_circle.fill.fore_color.rgb = needle_color
        center_circle.line.fill.background()
        
        # Score text
        score_text = slide.shapes.add_textbox(
            left + Inches(size/3), top + Inches(size/3), 
            Inches(size/3), Inches(size/3)
        )
        tf = score_text.text_frame
        tf.vertical_anchor = 1
        p = tf.paragraphs[0]
        p.text = str(risk_score)
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        
        # Labels: left = Critical (low score), right = Good (high score)
        label_positions = [
            ("Critical",  left - Inches(0.2),        top + Inches(size * 0.7)),
            ("High Risk", left + Inches(0.2),        top + Inches(size * 0.1)),
            ("Medium",    left + Inches(size * 0.7), top - Inches(0.3)),
            ("Good",      left + Inches(size * 0.8), top + Inches(size * 0.7)),
        ]
        
        for label_text, lx, ly in label_positions:
            label = slide.shapes.add_textbox(lx, ly, Inches(0.8), Inches(0.3))
            tf = label.text_frame
            p = tf.paragraphs[0]
            p.text = label_text
            p.font.size = Pt(9)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            
            # Color code the labels
            if label_text == "Critical":
                p.font.color.rgb = self.COLORS['critical']
            elif label_text == "High Risk":
                p.font.color.rgb = self.COLORS['high']
            elif label_text == "Medium":
                p.font.color.rgb = self.COLORS['medium']
            elif label_text == "Good":
                p.font.color.rgb = self.COLORS['green']
            else:
                p.font.color.rgb = self.COLORS['secondary']
        
        # Title above
        title = slide.shapes.add_textbox(left, top - Inches(0.4), Inches(size), Inches(0.3))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = "Riskometer"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['primary']
        p.alignment = PP_ALIGN.CENTER
    
    def _create_executive_summary_enhanced(self, prs, analysis, data):
        """Enhanced executive summary with riskometer and industry benchmark"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title bar
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(0.7)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.COLORS['primary']
        title_bar.line.fill.background()
        
        title = title_bar.text_frame
        title.text = f"{data['phase1_business']['company_name']} Security Intelligence Report"
        p = title.paragraphs[0]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        
        # --- LEFT SECTION: Risk Scores ---
        # Your Risk Score
        your_score_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1), Inches(2), Inches(1.2)
        )
        your_score_box.fill.solid()
        your_score_box.fill.fore_color.rgb = self.COLORS['accent']
        your_score_box.line.fill.background()
        
        self._add_wrapped_text(
            slide, "Your Risk Score", 
            Inches(0.5), Inches(1.05), Inches(2), Inches(0.4),
            font_size=14, bold=True, color=self.COLORS['white']
        )
        
        risk_score = data['risk_metrics']['risk_score']
        self._add_wrapped_text(
            slide, str(risk_score), 
            Inches(0.5), Inches(1.5), Inches(2), Inches(0.6),
            font_size=48, bold=True, color=self.COLORS['white']
        )
        
        # Industry Benchmark
        benchmark_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(2.3), Inches(2), Inches(1.2)
        )
        benchmark_box.fill.solid()
        benchmark_box.fill.fore_color.rgb = self.COLORS['secondary']
        benchmark_box.line.fill.background()
        
        self._add_wrapped_text(
            slide, "Industry Benchmark", 
            Inches(0.5), Inches(2.35), Inches(2), Inches(0.4),
            font_size=14, bold=True, color=self.COLORS['white']
        )
        
        benchmark_score = data['risk_metrics']['industry_benchmark']
        self._add_wrapped_text(
            slide, str(benchmark_score), 
            Inches(0.5), Inches(2.8), Inches(2), Inches(0.6),
            font_size=48, bold=True, color=self.COLORS['white']
        )
        
        # --- RISKOMETER ---
        self._create_riskometer(slide, risk_score, Inches(0.3), Inches(3.8))
        
        # --- TOP RIGHT: Executive Brief ---
        brief_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2.8), Inches(1), Inches(6.7), Inches(2.3)
        )
        brief_box.fill.solid()
        brief_box.fill.fore_color.rgb = self.COLORS['light_grey']
        brief_box.line.color.rgb = self.COLORS['accent']
        brief_box.line.width = Pt(2)
        
        # Title
        self._add_wrapped_text(
            slide, "Cybersecurity Executive Brief",
            Inches(2.95), Inches(1.15), Inches(6.4), Inches(0.35),
            font_size=16, bold=True, color=self.COLORS['primary']
        )
        
        # Executive summary text
        exec_text = analysis['comprehensive_analysis']['executive_summary']
        self._add_wrapped_text(
            slide, exec_text,
            Inches(3), Inches(1.6), Inches(6.3), Inches(1.5),
            font_size=12, color=self.COLORS['text']
        )
        
        # --- BOTTOM: Cyberattack Entry Points ---
        entry_header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(2.8), Inches(3.5), Inches(6.7), Inches(0.4)
        )
        entry_header.fill.solid()
        entry_header.fill.fore_color.rgb = self.COLORS['critical']
        entry_header.line.fill.background()
        
        self._add_wrapped_text(
            slide, "Cyberattack Entry Points Identified",
            Inches(2.8), Inches(3.52), Inches(6.7), Inches(0.35),
            font_size=14, bold=True, color=self.COLORS['white']
        )
        
        # Entry points (up to 3)
        entry_points = analysis.get('cyberattack_entry_points', [])[:3]
        y_pos = 4.0
        
        for idx, entry in enumerate(entry_points):
            # Entry point box
            entry_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(2.9), Inches(y_pos), Inches(6.5), Inches(0.85)
            )
            entry_box.fill.solid()
            entry_box.fill.fore_color.rgb = RGBColor(255, 250, 240)
            entry_box.line.color.rgb = self.COLORS['high']
            entry_box.line.width = Pt(1.5)
            
            # Title
            title_text = entry.get('title', 'Security Issue')
            if len(title_text) > 65:
                title_text = title_text[:62] + "..."
            
            self._add_wrapped_text(
                slide, title_text,
                Inches(3), Inches(y_pos + 0.05), Inches(6.3), Inches(0.3),
                font_size=11, bold=True, color=self.COLORS['critical']
            )
            
            # Description
            desc_text = entry.get('description', '')
            if len(desc_text) > 180:
                desc_text = desc_text[:177] + "..."
            
            self._add_wrapped_text(
                slide, desc_text,
                Inches(3), Inches(y_pos + 0.35), Inches(6.3), Inches(0.45),
                font_size=10, color=self.COLORS['text']
            )
            
            y_pos += 0.95
        
        # --- BOTTOM INFO ---
        # Perspective note
        perspective = slide.shapes.add_textbox(Inches(2.9), Inches(6.8), Inches(6.5), Inches(0.5))
        tf = perspective.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "An external perspective, emulating tactics employed by real-world attackers using passive reconnaissance to identify vulnerabilities."
        p.font.size = Pt(9)
        p.font.italic = True
        p.font.color.rgb = self.COLORS['dark_grey']
        
        # Last scanned
        scan_date = slide.shapes.add_textbox(Inches(8.5), Inches(7.2), Inches(1), Inches(0.2))
        tf = scan_date.text_frame
        p = tf.paragraphs[0]
        p.text = f"Last Scanned: {datetime.now().strftime('%B %Y')}"
        p.font.size = Pt(8)
        p.font.color.rgb = self.COLORS['dark_grey']
        p.alignment = PP_ALIGN.RIGHT
    
    def _create_vulnerability_analysis(self, prs, analysis, data):
        """Vulnerability analysis with charts"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(0.7)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.COLORS['primary']
        title_bar.line.fill.background()
        
        title = title_bar.text_frame
        title.text = "VULNERABILITY LANDSCAPE"
        p = title.paragraphs[0]
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        
        # CVE Distribution Chart
        chart_data = CategoryChartData()
        chart_data.categories = ['Critical', 'High', 'Medium', 'Low']
        chart_data.add_series('CVEs', (
            data['phase4_vulnerabilities']['critical_cves'],
            data['phase4_vulnerabilities']['high_cves'],
            data['phase4_vulnerabilities']['medium_cves'],
            data['phase4_vulnerabilities']['low_cves']
        ))
        
        x, y, cx, cy = Inches(0.5), Inches(1.2), Inches(4), Inches(3.8)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart
        
        chart.has_legend = False
        chart.has_title = True
        chart.chart_title.text_frame.text = "CVE Severity Distribution"
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(18)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
        
        # Vulnerability Summary Boxes
        summaries = [
            ("Critical", data['phase4_vulnerabilities']['critical_cves'], self.COLORS['critical']),
            ("High", data['phase4_vulnerabilities']['high_cves'], self.COLORS['high']),
            ("Medium", data['phase4_vulnerabilities']['medium_cves'], self.COLORS['medium']),
            ("Low", data['phase4_vulnerabilities']['low_cves'], self.COLORS['low'])
        ]
        
        y_start = 1.2
        for idx, (label, count, color) in enumerate(summaries):
            # Box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(5.3), Inches(y_start + idx * 1.45), Inches(4.2), Inches(1.3)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = color
            box.line.fill.background()
            
            # Label
            label_text = slide.shapes.add_textbox(
                Inches(5.5), Inches(y_start + idx * 1.45 + 0.15), Inches(3.8), Inches(0.3)
            )
            tf = label_text.text_frame
            p = tf.paragraphs[0]
            p.text = f"{label} Vulnerabilities"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.COLORS['white']
            
            # Count
            count_text = slide.shapes.add_textbox(
                Inches(5.5), Inches(y_start + idx * 1.45 + 0.5), Inches(3.8), Inches(0.4)
            )
            tf = count_text.text_frame
            p = tf.paragraphs[0]
            p.text = str(count)
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = self.COLORS['white']
            p.alignment = PP_ALIGN.CENTER
        
        # Risk Analysis Box
        analysis_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(6.5), Inches(9), Inches(0.8)
        )
        analysis_box.fill.solid()
        analysis_box.fill.fore_color.rgb = self.COLORS['light_grey']
        analysis_box.line.color.rgb = self.COLORS['info']
        analysis_box.line.width = Pt(2)
        
        risk_narrative = analysis['executive_summary'].get('risk_narrative', 'Comprehensive vulnerability assessment completed.')
        self._add_wrapped_text(
            slide, risk_narrative,
            Inches(0.7), Inches(6.6), Inches(8.6), Inches(0.6),
            font_size=12, color=self.COLORS['text']
        )
    
    def _create_findings_slide(self, prs, analysis, data):
        """Intelligence findings slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(0.7)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.COLORS['primary']
        title_bar.line.fill.background()
        
        title = title_bar.text_frame
        title.text = "KEY INTELLIGENCE FINDINGS"
        p = title.paragraphs[0]
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        
        # Findings
        findings = analysis.get('intelligent_findings', {})
        y_pos = 1.2
        
        for key, finding in list(findings.items())[:3]:
            # Finding box
            finding_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(y_pos), Inches(9), Inches(1.9)
            )
            finding_box.fill.solid()
            finding_box.fill.fore_color.rgb = self.COLORS['light_grey']
            
            severity = finding.get('severity', 'MEDIUM')
            if severity == 'CRITICAL':
                finding_box.line.color.rgb = self.COLORS['critical']
            elif severity == 'HIGH':
                finding_box.line.color.rgb = self.COLORS['high']
            elif severity == 'MEDIUM':
                finding_box.line.color.rgb = self.COLORS['medium']
            else:
                finding_box.line.color.rgb = self.COLORS['low']
            
            finding_box.line.width = Pt(3)
            
            # Severity badge
            badge = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(8.7), Inches(y_pos + 0.1), Inches(0.7), Inches(0.3)
            )
            badge.fill.solid()
            badge.fill.fore_color.rgb = finding_box.line.color.rgb
            badge.line.fill.background()
            
            badge_text = slide.shapes.add_textbox(
                Inches(8.7), Inches(y_pos + 0.1), Inches(0.7), Inches(0.3)
            )
            tf = badge_text.text_frame
            p = tf.paragraphs[0]
            p.text = severity
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = self.COLORS['white']
            p.alignment = PP_ALIGN.CENTER
            
            # Title
            title_text = finding.get('title', 'Finding')
            if len(title_text) > 90:
                title_text = title_text[:87] + "..."
            
            self._add_wrapped_text(
                slide, title_text,
                Inches(0.7), Inches(y_pos + 0.15), Inches(7.8), Inches(0.35),
                font_size=14, bold=True, color=self.COLORS['primary']
            )
            
            # Synthesis
            synthesis = finding.get('synthesis', '')
            if len(synthesis) > 300:
                synthesis = synthesis[:297] + "..."
            
            self._add_wrapped_text(
                slide, synthesis,
                Inches(0.7), Inches(y_pos + 0.55), Inches(8.4), Inches(0.6),
                font_size=11, color=self.COLORS['text']
            )
            
            # Business Impact
            impact = finding.get('business_impact', '')
            if len(impact) > 220:
                impact = impact[:217] + "..."
            
            impact_label = slide.shapes.add_textbox(
                Inches(0.7), Inches(y_pos + 1.2), Inches(1.5), Inches(0.3)
            )
            tf = impact_label.text_frame
            p = tf.paragraphs[0]
            p.text = "Business Impact:"
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = self.COLORS['high']
            
            self._add_wrapped_text(
                slide, impact,
                Inches(2.3), Inches(y_pos + 1.2), Inches(6.8), Inches(0.35),
                font_size=10, color=self.COLORS['text']
            )
            
            y_pos += 2.1
    
    def _create_roadmap(self, prs, analysis, data):
        """Strategic roadmap slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(0.7)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.COLORS['primary']
        title_bar.line.fill.background()
        
        title = title_bar.text_frame
        title.text = "STRATEGIC SECURITY ROADMAP"
        p = title.paragraphs[0]
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        
        # Immediate Priorities Section
        imm_header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.2), Inches(9), Inches(0.5)
        )
        imm_header.fill.solid()
        imm_header.fill.fore_color.rgb = self.COLORS['critical']
        imm_header.line.fill.background()
        
        self._add_wrapped_text(
            slide, "IMMEDIATE PRIORITIES",
            Inches(0.7), Inches(1.25), Inches(6), Inches(0.4),
            font_size=18, bold=True, color=self.COLORS['white']
        )
        
        self._add_wrapped_text(
            slide, "Next 30 Days",
            Inches(7.5), Inches(1.28), Inches(1.8), Inches(0.35),
            font_size=14, color=self.COLORS['white']
        )
        
        # Immediate items
        priorities = analysis['strategic_recommendations']['immediate_priorities']
        y_pos = 1.9
        
        for item in priorities:
            # Truncate if too long
            if len(item) > 180:
                item = item[:177] + "..."
            
            bullet_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.7), Inches(y_pos), Inches(8.6), Inches(0.6)
            )
            bullet_box.fill.solid()
            bullet_box.fill.fore_color.rgb = RGBColor(255, 240, 240)
            bullet_box.line.color.rgb = self.COLORS['high']
            bullet_box.line.width = Pt(1)
            
            self._add_wrapped_text(
                slide, f"✓ {item}",
                Inches(0.85), Inches(y_pos + 0.08), Inches(8.3), Inches(0.45),
                font_size=12, color=self.COLORS['text']
            )
            
            y_pos += 0.7
        
        # Strategic Initiatives Section
        strat_y = y_pos + 0.3
        strat_header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(strat_y), Inches(9), Inches(0.5)
        )
        strat_header.fill.solid()
        strat_header.fill.fore_color.rgb = self.COLORS['info']
        strat_header.line.fill.background()
        
        self._add_wrapped_text(
            slide, "STRATEGIC INITIATIVES",
            Inches(0.7), Inches(strat_y + 0.05), Inches(6), Inches(0.4),
            font_size=18, bold=True, color=self.COLORS['white']
        )
        
        self._add_wrapped_text(
            slide, "90+ Days",
            Inches(7.5), Inches(strat_y + 0.08), Inches(1.8), Inches(0.35),
            font_size=14, color=self.COLORS['white']
        )
        
        # Strategic items
        initiatives = analysis['strategic_recommendations']['strategic_initiatives']
        y_pos = strat_y + 0.7
        
        for item in initiatives:
            # Truncate if too long
            if len(item) > 180:
                item = item[:177] + "..."
            
            bullet_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.7), Inches(y_pos), Inches(8.6), Inches(0.6)
            )
            bullet_box.fill.solid()
            bullet_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
            bullet_box.line.color.rgb = self.COLORS['info']
            bullet_box.line.width = Pt(1)
            
            self._add_wrapped_text(
                slide, f"✓ {item}",
                Inches(0.85), Inches(y_pos + 0.08), Inches(8.3), Inches(0.45),
                font_size=12, color=self.COLORS['text']
            )
            
            y_pos += 0.7
    
    def _get_risk_color(self, score):
        """Get color based on risk score"""
        if score >= 80:
            return self.COLORS['low']
        elif score >= 60:
            return self.COLORS['medium']
        elif score >= 40:
            return self.COLORS['high']
        else:
            return self.COLORS['critical']
    
    def _extract_nested(self, data, keys, default=None):
        """Safely extract nested dictionary values"""
        try:
            result = data
            for key in keys:
                result = result[key]
            return result
        except (KeyError, TypeError):
            return default
    
    def generate(self, phase1, phase2, phase3, phase4, phase5, output_path=None):
        """
        Main generation - AI analyzes ALL 5 phases comprehensively
        """
        print("\n" + "="*70)
        print("🧠 ENHANCED INTELLIGENT PROFESSIONAL REPORT GENERATOR")
        print("="*70)
        
        # AI comprehensive analysis
        ai_analysis, comprehensive_data = self.comprehensive_ai_analysis(
            phase1, phase2, phase3, phase4, phase5
        )
        
        # Create presentation
        if not output_path:
            output_path = os.getcwd()
        
        filepath = self.create_intelligent_presentation(
            ai_analysis, comprehensive_data, output_path
        )
        
        print("\n" + "="*70)
        print(f"✅ ENHANCED INTELLIGENT REPORT COMPLETE: {filepath}")
        print(f"   ✅ Risk Score & Industry Benchmark")
        print(f"   ✅ Riskometer Visualization")
        print(f"   ✅ Cyberattack Entry Points")
        print(f"   ✅ Improved Text Wrapping")
        print(f"   ✅ Professional Layout")
        print(f"   AI: {'COMPREHENSIVE 5-PHASE ANALYSIS' if self.use_ai else 'Template'}")
        print("="*70 + "\n")
        
        return filepath


if __name__ == "__main__":
    print(__doc__)